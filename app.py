import streamlit as st
from google import genai
from google.genai import types
import streamlit.components.v1 as components
from recipe_scrapers import scrape_me
import json
from io import BytesIO
import time
import requests
from bs4 import BeautifulSoup

# Importaciones seguras para documentos opcionales
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

st.set_page_config(page_title="FaceFoodChef Pro - Multimodal", layout="centered", page_icon="🍳")

# Estilos visuales FaceFoodChef (Modo Oscuro Oficial)
st.markdown("""
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500&family=JetBrains+Mono:wght@400;700&family=Montserrat:wght@700;800&display=swap" rel="stylesheet">
    <style>
    .stApp { background-color: #2C2F33; font-family: 'Inter', sans-serif; color: #E2E8F0; }
    .stTextArea textarea, .stFileUploader {
        background-color: #36393F; color: #E2E8F0;
        border-radius: 8px; border: 1px solid #4F545C; font-size: 15px;
        font-family: 'Inter', sans-serif;
    }
    h1, h2, h3, .st-emotion-cache-10trblm { color: #FFFFFF; font-family: 'Montserrat', sans-serif !important; font-weight: 700; }
    p, span, div { font-family: 'Inter', sans-serif; }
    code { font-family: 'JetBrains Mono', monospace !important; }
    </style>
""", unsafe_allow_html=True)

st.markdown("<h1 style='text-align: center; color: #FFFFFF; border-bottom: 2px solid #EF4444; padding-bottom: 10px;'>🍳 FaceFoodChef Pro</h1>", unsafe_allow_html=True)
st.markdown("<p style='text-align: center; color: #E2E8F0;'>Ingeniería de procesos culinarios avanzada. Sube PDFs, Word, PPT, Imágenes o pega URLs y textos.</p>", unsafe_allow_html=True)

st.sidebar.header("⚙️ Panel de Control")
API_KEY = st.sidebar.text_input("API Key de Google Gemini:", type="password", help="Consíguela gratis en aistudio.google.com")
st.sidebar.markdown("---")
st.sidebar.markdown("### 📂 Formatos Compatibles\n- **Web:** Cualquier URL (Blogs, Periódicos, etc.).\n- **Documentos:** PDF, Word (.docx), PPT (.pptx).\n- **Imágenes:** JPG, PNG, WEBP.\n- **Texto:** Directo o apuntes.")

tab1, tab2, tab3 = st.tabs(["🌐 URL o Texto", "📁 Subir Archivo (PDF, Word, PPT)", "🖼️ Subir Imagen de Receta"])

receta_texto_input = ""
archivo_multimodal = None
tipo_multimodal = None

with tab1:
    entrada_usuario = st.text_area(
        "📝 Pega la URL de una receta web o escribe los pasos manualmente:", 
        height=100, 
        placeholder="Ej: https://okdiario.com/... o escribe tu receta..."
    )
    if entrada_usuario:
        receta_texto_input = entrada_usuario

with tab2:
    archivo_doc = st.file_uploader("📂 Sube un documento con tu receta:", type=["pdf", "docx", "pptx", "txt"])
    if archivo_doc:
        ext = archivo_doc.name.split('.')[-1].lower()
        if ext == "txt":
            receta_texto_input = archivo_doc.getvalue().decode("utf-8")
        elif ext == "docx":
            if HAS_DOCX:
                doc = docx.Document(BytesIO(archivo_doc.getvalue()))
                receta_texto_input = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            else:
                st.error("⚠️ La librería 'python-docx' no está instalada.")
        elif ext == "pptx":
            if HAS_PPTX:
                prs = Presentation(BytesIO(archivo_doc.getvalue()))
                texto_slides = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                texto_slides.append(paragraph.text)
                receta_texto_input = "\n".join(texto_slides)
            else:
                st.error("⚠️ La librería 'python-pptx' no está instalada.")
        elif ext == "pdf":
            archivo_multimodal = archivo_doc.getvalue()
            tipo_multimodal = "application/pdf"

with tab3:
    archivo_img = st.file_uploader("📸 Sube una foto de receta o captura de pantalla:", type=["jpg", "jpeg", "png", "webp"])
    if archivo_img:
        archivo_multimodal = archivo_img.getvalue()
        tipo_multimodal = archivo_img.type

def extraer_texto_de_url(url):
    url = url.strip()
    try:
        scraper = scrape_me(url)
        return f"Receta extraída de {url}:\nIngredientes: {', '.join(scraper.ingredients())}\nInstrucciones:\n{'\n'.join(scraper.instructions())}"
    except Exception:
        try:
            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'}
            response = requests.get(url, headers=headers, timeout=10)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            for element in soup(["script", "style", "nav", "footer", "header", "aside"]):
                element.decompose()
            texto_limpio = soup.get_text(separator='\n', strip=True)
            return f"Contenido extraído de la URL ({url}):\n{texto_limpio}"
        except Exception as e_general:
            raise Exception(f"No se pudo leer la URL. Comprueba que sea accesible. Detalle: {e_general}")

def generar_html_dashboard(ingredientes, pasos_previos, bloques_proceso, texto_voz):
    html_ing = """
    <div style="background: #36393F; border: 1px solid #4F545C; border-radius: 12px; padding: 22px; margin-bottom: 20px;">
        <h3 style="color: #FFF; font-family: 'Montserrat', sans-serif; margin-top: 0; font-size: 17px; display: flex; align-items: center; gap: 8px;">
            <span style="color: #EF4444;">🛒</span> 1. Despensa e Ingredientes
        </h3>
        <div style="display: flex; flex-wrap: wrap; gap: 8px; margin-top: 14px;">
    """
    for ing in ingredientes:
        html_ing += f"<span style='background: #2C2F33; color: #E2E8F0; padding: 6px 14px; border-radius: 6px; font-size: 13px; border: 1px solid #4F545C;'>{ing}</span>"
    html_ing += "</div></div>"

    html_prev = """
    <div style="background: #36393F; border: 1px solid #4F545C; border-radius: 12px; padding: 22px; margin-bottom: 24px;">
        <h3 style="color: #FFF; font-family: 'Montserrat', sans-serif; margin-top: 0; font-size: 17px; display: flex; align-items: center; gap: 8px;">
            <span style="color: #EF4444;">🔪</span> 2. Mise en Place (Preparación Previa)
        </h3>
        <ul style='margin: 10px 0 0 0; padding-left: 20px; color: #E2E8F0; font-size: 14px; line-height: 1.7;'>
    """
    for prep in pasos_previos:
        html_prev += f"<li style='margin-bottom: 6px;'>{prep}</li>"
    html_prev += "</ul></div>"

    html_diagrama = """
    <div style="max-width: 100%; margin: auto;">
        <h3 style="color: #FFF; font-family: 'Montserrat', sans-serif; font-size: 18px; margin-bottom: 20px; display: flex; align-items: center; gap: 8px;">
            <span style="color: #EF4444;">📊</span> 3. Diagrama de Flujo Interactivo
        </h3>
    """
    
    for i, bloque in enumerate(bloques_proceso):
        tipo = bloque.get("tipo", "secuencial")
        duracion_min = bloque.get("duracion_minutos", 5)
        
        if tipo == "paralelo":
            ramas = bloque.get("ramas", [])
            html_diagrama += '<div style="display: flex; gap: 14px; margin-bottom: 18px; flex-wrap: wrap;">'
            for idx, rama in enumerate(ramas):
                nombre_rama = rama.get("nombre", f"Rama {idx+1}").upper()
                accion = rama.get("accion", "")
                tiempo = rama.get("tiempo", "")
                temp = rama.get("temperatura", "")
                dur_rama = rama.get("duracion_minutos", 5)
                timer_id = f"timer_par_{i}_{idx}"
                
                html_diagrama += f"""
                <div style="flex: 1; min-width: 270px; background: #2C2F33; border: 1px solid #4F545C; border-left: 4px solid #EF4444; border-radius: 8px; padding: 20px;">
                    <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 10px;">
                        <span style="font-size: 11px; font-weight: 700; color: #FFF; background: #EF4444; padding: 4px 8px; border-radius: 4px;">🔀 PARALELO: {nombre_rama}</span>
                    </div>
                    <div style="font-size: 15px; color: #E2E8F0; margin-bottom: 14px; line-height: 1.5;">{accion}</div>
                    <div style="display: flex; justify-content: space-between; align-items: center; margin-top: 14px; background: #36393F; padding: 10px 14px; border-radius: 6px; font-family: 'JetBrains Mono', monospace;">
                        <div style="font-size: 13px; color: #E2E8F0;">⏱️ <span id="{timer_id}" style="font-weight: bold; color: #EF4444;">{tiempo}</span> | 🌡️ {temp}</div>
                        <button onclick="iniciarTemporizador('{timer_id}', {dur_rama})" style="background: #EF4444; color: white; border: none; padding: 6px 12px; border-radius: 4px; cursor: pointer; font-size: 12px; font-family: 'Montserrat', sans-serif; font-weight: 700;">▶️ Iniciar</button>
                    </div>
                </div>
                """
            html_diagrama += '</div>'
        else:
            es_convergencia = tipo == "convergencia"
            bg = "#2C2F33" if es_convergencia else "#36393F"
            icono = "🔗" if es_convergencia else "🔥"
            etiqueta = "CONVERGENCIA" if es_convergencia else f"BLOQUE SECUENCIAL {i+1}"
            
            timer_id = f"timer_seq_{i}"
            html_diagrama += f"""
            <div style="background: {bg}; border: 1px solid #4F545C; border-left: 4px solid #EF4444; border-radius: 8px; padding: 20px; margin-bottom: 18px;">
                <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 10px;">
                    <span style="font-size: 11px; font-weight: 700; color: #FFF; background: #EF4444; padding: 4px 8px; border-radius: 4px;">{icono} {etiqueta}</span>
                </div>
                <div style="font-size: 15px; color: #E2E8F0; margin-bottom: 14px; line-height: 1.5;">{bloque.get('accion')}</div>
                <div style="display: flex; justify-content: space-between; align-items: center; background: #2C2F33; padding: 10px 14px; border-radius: 6px; font-family: 'JetBrains Mono', monospace;">
                    <div style="font-size: 13px; color: #E2E8F0;">⏱️ <span id="{timer_id}" style="font-weight: bold; color: #EF4444;">{bloque.get('tiempo')}</span> | 🌡️ {bloque.get('temperatura')}</div>
                    <button onclick="iniciarTemporizador('{timer_id}', {duracion_min})" style="background: #EF4444; color: white; border: none; padding: 6px 12px; border-radius: 4px; cursor: pointer; font-size: 12px; font-family: 'Montserrat', sans-serif; font-weight: 700;">▶️ Iniciar</button>
                </div>
            </div>
            """
            
        if i < len(bloques_proceso) - 1:
            html_diagrama += '<div style="text-align: center; color: #EF4444; font-size: 20px; margin: -6px 0 6px 0; font-weight: bold;">↓</div>'
    
    html_diagrama += "</div>"

    documento_completo = f"""
    <!DOCTYPE html>
    <html lang="es">
    <head>
        <meta charset="utf-8">
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500&family=JetBrains+Mono:wght@400;700&family=Montserrat:wght@700;800&display=swap" rel="stylesheet">
        <style>
            body {{ background-color: #2C2F33; color: #E2E8F0; font-family: 'Inter', sans-serif; padding: 16px; margin: 0; }}
            .container-hub {{ display: grid; grid-template-columns: 1fr; gap: 20px; max-width: 900px; margin: auto; }}
            .widget-box {{ background: #36393F; border: 1px solid #4F545C; border-radius: 12px; padding: 18px; text-align: center; }}
            .btn-control {{ background-color: #EF4444; color: white; border: none; padding: 10px 20px; font-size: 13px; font-family: 'Montserrat', sans-serif; font-weight: 700; border-radius: 6px; cursor: pointer; margin: 4px; }}
            .btn-control:hover {{ opacity: 0.9; }}
            .btn-stop {{ background-color: #4F545C; }}
        </style>
    </head>
    <body>
        <div class="container-hub">
            <div class="widget-box">
                <p style="color: #E2E8F0; font-size: 13px; margin: 0 0 10px 0; font-family: 'Montserrat', sans-serif; font-weight: 700;">🎧 Tu Música Favorita (Spotify)</p>
                <iframe style="border-radius:8px" src="https://open.spotify.com/embed/playlist/37i9dQZF1DXcBWIGoYBM5M?utm_source=generator&theme=0" width="100%" height="80" frameBorder="0" allowfullscreen="" allow="autoplay; clipboard-write; encrypted-media; fullscreen; picture-in-picture" loading="lazy"></iframe>
            </div>
            <div class="widget-box">
                <p style="color: #E2E8F0; font-size: 13px; margin: 0 0 10px 0; font-family: 'Montserrat', sans-serif; font-weight: 700;">🎙️ Asistente de Voz Integrado</p>
                <button class="btn-control" onclick="reproducir()">▶️ Escuchar Guía</button>
                <button class="btn-control btn-stop" onclick="detener()">⏹️ Silenciar</button>
            </div>
            {html_ing}
            {html_prev}
            {html_diagrama}
        </div>
        <script>
            const textoVoz = "{texto_voz.replace('"', '').replace(chr(10), '. ')}";
            function reproducir() {{
                if ('speechSynthesis' in window) {{
                    window.speechSynthesis.cancel();
                    const msg = new SpeechSynthesisUtterance(textoVoz);
                    msg.lang = 'es-ES'; msg.rate = 0.95;
                    window.speechSynthesis.speak(msg);
                }}
            }}
            function detener() {{
                if ('speechSynthesis' in window) {{ window.speechSynthesis.cancel(); }}
            }}
            function iniciarTemporizador(elementId, minutos) {{
                const elemento = document.getElementById(elementId);
                let segundosRestantes = minutos * 60;
                if (window[elementId + "_interval"]) {{ clearInterval(window[elementId + "_interval"]); }}
                window[elementId + "_interval"] = setInterval(() => {{
                    if (segundosRestantes <= 0) {{
                        clearInterval(window[elementId + "_interval"]);
                        elemento.innerText = "¡TIEMPO CUMPLIDO! 🎉";
                        sonarAlerta();
                    }} else {{
                        segundosRestantes--;
                        const m = Math.floor(segundosRestantes / 60);
                        const s = segundosRestantes % 60;
                        elemento.innerText = `${{m}}m ${{s < 10 ? '0' : ''}}${{s}}s`;
                    }}
                }}, 1000);
            }}
            function sonarAlerta() {{
                const audioCtx = new (window.AudioContext || window.webkitAudioContext)();
                const oscillator = audioCtx.createOscillator();
                const gainNode = audioCtx.createGain();
                oscillator.type = 'sine';
                oscillator.frequency.setValueAtTime(587.33, audioCtx.currentTime);
                gainNode.gain.setValueAtTime(0.2, audioCtx.currentTime);
                oscillator.connect(gainNode);
                gainNode.connect(audioCtx.destination);
                oscillator.start();
                setTimeout(() => {{ oscillator.stop(); }}, 1200);
            }}
        </script>
    </body>
    </html>
    """
    return documento_completo

procesar_accion = False
contenido_ia = None

if API_KEY:
    if receta_texto_input.strip().startswith("http://") or receta_texto_input.strip().startswith("https://"):
        try:
            with st.spinner("🌐 Extrayendo datos de la web (con fallback universal)..."):
                contenido_ia = extraer_texto_de_url(receta_texto_input.strip())
                procesar_accion = True
        except Exception as e:
            st.error(f"{e}")
    elif receta_texto_input.strip():
        contenido_ia = receta_texto_input
        procesar_accion = True
    elif archivo_multimodal:
        procesar_accion = True

if procesar_accion and API_KEY:
    try:
        client = genai.Client(api_key=API_KEY)
        
        prompt_sistema = """
        Eres un chef ejecutivo e ingeniero de procesos culinarios. Analiza la fuente aportada (texto plano, artículo web extraído, documento o imagen) y estructúrala detectando tareas secuenciales, en paralelo (columnas simultáneas) y convergencias finales.
        
        Devuélvela estrictamente en formato JSON válido con esta estructura exacta, asegurándote de incluir el campo 'duracion_minutos' (número entero con los minutos estimados para cada tarea, necesario para los temporizadores):
        
        REGLAS DE ORO:
        1. 'ingredientes': Lista de ingredientes con cantidades exactas.
        2. 'pasos_previos': Lista con la preparación previa (mise en place).
        3. 'bloques_proceso': Lista de objetos con los siguientes tipos:
           - TIPO 1: {"tipo": "secuencial", "accion": "...", "tiempo": "5 min", "duracion_minutos": 5, "temperatura": "..."}
           - TIPO 2: {"tipo": "paralelo", "ramas": [{"nombre": "...", "accion": "...", "tiempo": "10 min", "duracion_minutos": 10, "temperatura": "..."}, {"nombre": "...", "accion": "...", "tiempo": "8 min", "duracion_minutos": 8, "temperatura": "..."}]}
           - TIPO 3: {"tipo": "convergencia", "accion": "...", "tiempo": "2 min", "duracion_minutos": 2, "temperatura": "..."}
        4. 'texto_voz': Resumen estructurado para lectura por voz.

        Estructura JSON obligatoria:
        {
          "ingredientes": ["..."],
          "pasos_previos": ["..."],
          "bloques_proceso": [
            {"tipo": "secuencial", "accion": "🔥 ...", "tiempo": "5 min", "duracion_minutos": 5, "temperatura": "Fuego medio"}
          ],
          "texto_voz": "..."
        }
        """

        contents_payload = [prompt_sistema]
        if archivo_multimodal:
            contents_payload.append(types.Part.from_bytes(data=archivo_multimodal, mime_type=tipo_multimodal))
            contents_payload.append("Analiza esta fuente adjunta y extrae la receta completa.")
        else:
            contents_payload.append(f"Información a procesar:\n{contenido_ia}")

        response = None
        max_intentos = 3
        
        for intento in range(max_intentos):
            try:
                with st.spinner(f"⚙️ Procesando con IA (Intento {intento+1}/{max_intentos})..."):
                    response = client.models.generate_content(
                        model='gemini-2.5-flash',
                        contents=contents_payload,
                        config=types.GenerateContentConfig(
                            response_mime_type="application/json",
                            temperature=0.2
                        ),
                    )
                break
            except Exception as api_err:
                err_str = str(api_err)
                if ("503" in err_str or "UNAVAILABLE" in err_str) and intento < max_intentos - 1:
                    time.sleep(3 * (intento + 1))
                    continue
                else:
                    raise api_err

        if response:
            texto_respuesta = response.text.strip()
            if texto_respuesta.startswith("```json"):
                texto_respuesta = texto_respuesta[7:]
            if texto_respuesta.endswith("```"):
                texto_respuesta = texto_respuesta[:-3]
            
            datos = json.loads(texto_respuesta.strip())
            
            html_final = generar_html_dashboard(
                datos.get("ingredientes", []),
                datos.get("pasos_previos", []),
                datos.get("bloques_proceso", []),
                datos.get("texto_voz", "")
            )
            
            st.download_button(
                label="📥 Descargar Dashboard Pro Completo (HTML)",
                data=html_final,
                file_name="facefoodchef_pro_hub.html",
                mime="text/html"
            )
            
            components.html(html_final, height=1050, scrolling=True)
            
    except Exception as e:
        st.error(f"Error procesando con la IA: {e}")
elif not API_KEY:
    st.info("👈 Introduce tu API Key gratuita de Google AI Studio en la barra lateral para comenzar.")
else:
    st.info("👆 Selecciona una de las pestañas superiores para introducir una URL, subir un documento (PDF, Word, PPT) o una foto de tu receta.")
